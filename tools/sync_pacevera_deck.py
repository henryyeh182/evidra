from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SRC = "docs/Fitness-MCP-Product-Vision.pptx"
OUT = "docs/Fitness-MCP-Product-Vision.pptx"

PAPER = "F0EEE6"
SURFACE = "FAF9F5"
SURFACE2 = "E9E3D6"
LINE = "DAD2C3"
INK = "1F1E1D"
MUTED = "68645C"
MINT = "157554"
MINT_SOFT = "DCEDE4"
NAVY = "111B30"
NAVY_CARD = "1E2A40"
NAVY_LINE = "34445E"
LIME = "49D993"
AMBER = "ECAE4D"
RED = "EC7972"


def rgb(value):
    return RGBColor.from_string(value)


def color_of(obj):
    try:
        if obj.type is None:
            return None
        return str(obj.rgb)
    except (AttributeError, ValueError, TypeError):
        return None


def set_color(obj, value):
    try:
        obj.rgb = rgb(value)
    except (AttributeError, ValueError, TypeError):
        pass


def set_east_asian_font(run, family="Microsoft JhengHei"):
    """Set CJK font slots too; setting only a:latin drops Traditional Chinese in some viewers."""
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        node = r_pr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            r_pr.append(node)
        node.set("typeface", family)


def replace_text(slide, replacements):
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text:
            continue
        text = shape.text
        for old, new in replacements.items():
            text = text.replace(old, new)
        if text != shape.text:
            shape.text = text


def style_text(shape, dark):
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            set_east_asian_font(run)
            old = color_of(run.font.color)
            mapped = {
                "6E7687": MUTED,
                "A8AEBC": "9DA89F",
                "CFD4DF": "C9D0C8",
                "2C313D": INK if not dark else NAVY_LINE,
                "10131A": INK if not dark else PAPER,
                "353A46": INK if not dark else SURFACE,
                "8A93A6": MUTED if not dark else "9FAEBE",
                "3A4050": LINE if not dark else NAVY_LINE,
                "9EEBD8": MINT if not dark else LIME,
                "1FB6A6": MINT if not dark else LIME,
                "F96167": MINT if not dark else LIME,
                "FFE3E4": MINT_SOFT if not dark else "FFE0C4",
                "FFB3B6": MINT if not dark else AMBER,
                "FFFFFF": INK if not dark else SURFACE,
            }.get(old)
            if mapped:
                set_color(run.font.color, mapped)
            elif old is None:
                set_color(run.font.color, SURFACE if dark else INK)


def style_shape(shape, dark):
    try:
        fill_color = color_of(shape.fill.fore_color)
        fill_map = {
            "FFFFFF": NAVY if dark else PAPER,
            "F3F4F7": SURFACE2,
            "10131A": NAVY,
            "1B1F2B": NAVY_CARD,
            "1FB6A6": MINT,
            "F96167": MINT,
            "4A5568": NAVY_LINE,
            "2C3242": NAVY_LINE,
            "8A93A6": MUTED,
        }
        if fill_color in fill_map:
            set_color(shape.fill.fore_color, fill_map[fill_color])
    except (AttributeError, ValueError, TypeError):
        pass
    style_text(shape, dark)


def add_text(slide, text, x, y, w, h, size, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    set_east_asian_font(run)
    run.font.size = Pt(size)
    run.font.bold = bold
    set_color(run.font.color, color)
    return box


def add_flow_box(slide, x, w, fill, title, detail, title_color=SURFACE, detail_color="C9D0C8", accent=None, central=False):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.55), Inches(w), Inches(1.35))
    box.fill.solid()
    set_color(box.fill.fore_color, fill)
    box.line.fill.background()
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(.18)
    tf.margin_right = Inches(.18)
    tf.margin_top = Inches(.12)
    tf.margin_bottom = Inches(.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = "Arial"
    set_east_asian_font(r)
    r.font.size = Pt(23 if central else 17)
    r.font.bold = True
    set_color(r.font.color, accent or title_color)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(5)
    r2 = p2.add_run()
    r2.text = detail
    r2.font.name = "Arial"
    set_east_asian_font(r2)
    r2.font.size = Pt(10.5)
    set_color(r2.font.color, detail_color)
    return box


def rebuild_slide_11(prs):
    slide = prs.slides[10]
    slide.background.fill.solid()
    set_color(slide.background.fill.fore_color, PAPER)
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)

    add_text(slide, "PART 2 — CORE FUNCTION", .78, .48, 3.6, .25, 10, MUTED, True)
    add_text(slide, "AI can coach. Pacevera decides.", .78, .86, 8.9, .62, 30, INK, True)
    add_text(slide, "One natural-language question becomes an executable training change through evidence and deterministic rules.", .78, 1.55, 10.6, .32, 13, MUTED)
    add_text(slide, "ONE DECISION LOOP", .78, 2.16, 2.3, .2, 9, MINT, True)

    add_flow_box(slide, .78, 2.2, SURFACE2, "1  AI coach", "understand intent · call tools", title_color=INK, detail_color=MUTED, accent=INK, central=False)
    add_text(slide, "→", 3.06, 2.96, .35, .35, 22, MINT, True, PP_ALIGN.CENTER)
    add_flow_box(slide, 3.43, 2.55, NAVY, "2  Pacevera", "calculate · decide · trace", accent=LIME, central=True)
    add_text(slide, "→", 6.08, 2.96, .35, .35, 22, MINT, True, PP_ALIGN.CENTER)
    add_flow_box(slide, 6.45, 2.35, SURFACE2, "3  Evidence + rules", "recovery · fatigue · load · plan", title_color=INK, detail_color=MUTED, accent=INK)
    add_text(slide, "→", 8.91, 2.96, .35, .35, 22, MINT, True, PP_ALIGN.CENTER)
    add_flow_box(slide, 9.28, 3.25, SURFACE2, "4  Decision + reason", "keep · adjust · defer · change for today", title_color=INK, detail_color=MUTED, accent=INK, central=False)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(.78), Inches(4.38), Inches(11.75), Inches(1.38))
    panel.fill.solid(); set_color(panel.fill.fore_color, SURFACE2); panel.line.fill.background()
    add_text(slide, "PACEVERA RETURNS", 1.02, 4.62, 2.3, .25, 10, MINT, True)
    add_text(slide, "A change you can act on — with the evidence and limits attached.", 1.02, 4.94, 4.8, .42, 15, INK, True)
    add_text(slide, "planned → today", 6.35, 4.62, 1.2, .22, 9, MUTED, True)
    add_text(slide, "reason", 7.88, 4.62, 1.25, .22, 9, MUTED, True)
    add_text(slide, "missing", 9.55, 4.62, 1.15, .22, 9, MUTED, True)
    add_text(slide, "confidence", 11.0, 4.62, 1.2, .22, 9, MUTED, True)
    add_text(slide, "Tempo → Moderate", 6.35, 4.98, 1.25, .42, 10, INK, True)
    add_text(slide, "readiness below threshold", 7.88, 4.98, 1.35, .42, 9, INK)
    add_text(slide, "sleep", 9.55, 4.98, .9, .3, 10, INK)
    add_text(slide, "medium", 11.0, 4.98, .9, .3, 10, MINT, True)
    add_text(slide, "11", 12.2, 6.72, .35, .18, 9, MUTED, False, PP_ALIGN.RIGHT)


def main():
    prs = Presentation(SRC)
    slide_replacements = {
        1: {
            "讓你熟悉的 AI 根據連續 Evidence 做出可重現的訓練決策 —— 健康歷史仍由你控制。": "把 Evidence 轉成今天可執行、可追溯的訓練變更。健康歷史仍由你控制。",
            "產品終局 · v0.4.2 → v0.5.0 · Phase 0–4": "Public preview · Desktop MCPB · v0.5.0 → local private engine",
        },
        2: {
            "先看產品終局,再看已交付能力,最後用真實情境與 roadmap 驗收。": "先看今天的決策，再看 evidence chain 與 privacy boundary，最後看產品驗證與 roadmap。",
            "我們蓋的不是 App,是一個確定性的訓練決策引擎": "從原定課表到今天要做什麼",
            "目前 6 tools、下一版 10 tools,以及版本化的 Engine 與 Rule Packages": "證據、理由、缺失與信心都可檢查",
            "六個真實問題,以及從 v0.5.0 到 private engine 的交付順序": "AI coach 可替換，健康歷史留在你的控制邊界",
        },
        3: {"身體是連續的,資料卻是被切開的": "一般 AI 知道運動知識，卻不知道連續的你"},
        5: {"我們要蓋的,是中間那一層": "我們要蓋的是 AI coach 與今天決策之間的一層"},
        6: {"使用者要連續的決策能力,不要再交出健康歷史": "Your AI can coach. Pacevera helps it decide."},
        7: {"關鍵設計:資料由使用者控制": "Evidence 留在你的控制邊界"},
        8: {"輸出具體長這樣": "輸出不是分數，是今天的 change"},
        10: {"我們給的是決策,不是建議": "我們給的是決策，不是建議"},
        12: {"決策引擎不是一個模型": "先對齊 Evidence，再用規則做決策"},
        13: {"Planning Engine 是整個產品最大的價值": "從計畫到今天：改動要能執行"},
        17: {"各家指標對齊,廠商算好的分數直接收": "把不同來源對齊成同一套語彙"},
        18: {"Connector 跑在你的環境,不在我們的": "Connector 跑在你的環境"},
        19: {"v0.5.0 規劃的十個 public tools": "v0.5.0：把一個 decision loop 做完整"},
        21: {"三種部署,三句不同的隱私承諾": "同一顆 engine，三種部署形態"},
        22: {"六個使用者真的會問的問題": "六種使用者真的會問的問題"},
        29: {"六題的共通點:先有證據,才談決策": "共通點：先有 Evidence，才談 Decision"},
        30: {"從已發行版本倒推下一階段": "先用產品頁驗證，再進 private engine"},
        31: {
            "Plaid 之於金融,Stripe 之於支付": "Pacevera 是 personal fitness decision layer",
            "我們要的是運動領域的決策基礎層:證據、規則與可追溯的 from → to。": "我們要的是運動領域的決策基礎層:證據、規則與可追溯的訓練變更。",
        },
    }
    for number, replacements in slide_replacements.items():
        replace_text(prs.slides[number - 1], replacements)

    global_replacements = {
        "Fitness MCP": "Pacevera",
        "Fitness Decision Engine": "Fitness Decision Engine",
        "Pacevera\nFitness Decision Engine": "Pacevera\nFitness Decision Engine",
        "仲裁": "判斷",
        "Decision from → to · reason · missing signals · versioned trace": "Planned session · Today’s session · reason · missing signals · versioned trace",
        "並寫成 from → to。": "並寫成「原定課表到今天執行」。",
        "from → to": "planned session → today’s session",
        "from:": "planned_session:",
        "to:": "today_session:",
    }
    for slide in prs.slides:
        replace_text(slide, global_replacements)
        dark = color_of(slide.background.fill.fore_color) in {"10131A", "1B1F2B"}
        set_color(slide.background.fill.fore_color, NAVY if dark else PAPER)
        for shape in slide.shapes:
            style_shape(shape, dark)

    # Page 6: highlight the user's desired outcome, not the thing they reject.
    page6 = prs.slides[5]
    for index, fill in ((4, MINT), (7, NAVY_CARD)):
        shape = page6.shapes[index]
        shape.fill.solid()
        set_color(shape.fill.fore_color, fill)
    for index in (5, 6):
        shape = page6.shapes[index]
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                set_color(run.font.color, SURFACE)
    for index in (8, 9, 10):
        shape = page6.shapes[index]
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                set_color(run.font.color, SURFACE)

    # Contrast corrections: light text on navy, dark text on beige.
    def recolor_text(slide, indices, color):
        for index in indices:
            shape = slide.shapes[index]
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_color(run.font.color, color)

    page10 = prs.slides[9]
    recolor_text(page10, (6,), INK)
    recolor_text(page10, (9,), SURFACE)

    page17 = prs.slides[16]
    recolor_text(page17, (6, 10), INK)
    recolor_text(page17, (14,), SURFACE)

    page31 = prs.slides[30]
    recolor_text(page31, (5, 6, 8, 9), INK)
    recolor_text(page31, (12,), SURFACE)

    # Page 7 is the privacy-boundary summary. Keep both payload rows readable
    # against the homepage palette after the global colour migration.
    page7 = prs.slides[6]
    for shape in page7.shapes:
        if not hasattr(shape, "text"):
            continue
        if shape.text.startswith("Raw Evidence"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_color(run.font.color, INK)
        elif shape.text.startswith("Decision from"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_color(run.font.color, SURFACE)
    rebuild_slide_11(prs)
    prs.save(OUT)


if __name__ == "__main__":
    main()

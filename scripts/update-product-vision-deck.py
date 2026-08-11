#!/usr/bin/env python3
"""Align the existing Pacevera product-vision deck with the canonical roadmap.

The script intentionally updates text in place and preserves the existing slide
geometry, theme, images, and object ordering.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


REPLACEMENTS: dict[int, dict[str, str]] = {
    1: {
        "讓任何 AI 都能做出可重現的訓練決策 —— 而不必看見你完整的原始健康資料。":
            "讓你熟悉的 AI 根據連續 Evidence 做出可重現的訓練決策 —— 健康歷史仍由你控制。",
        "核心概念 · 核心功能 · 應用情境":
            "產品終局 · v0.4.2 → v0.5.0 · Phase 0–4",
    },
    2: {
        "先講清楚我們在蓋什麼,再講它能做什麼,最後用真實情境驗收。":
            "先看產品終局,再看已交付能力,最後用真實情境與 roadmap 驗收。",
        "六個 tool,以及背後那個不呼叫模型的決策引擎":
            "目前 6 tools、下一版 10 tools,以及版本化的 Engine 與 Rule Packages",
        "六個使用者真的會問的問題,以及每一題打到哪幾層":
            "六個真實問題,以及從 v0.5.0 到 private engine 的交付順序",
    },
    6: {
        "使用者要教練能力,不要交出原始資料":
            "使用者要連續的決策能力,不要再交出健康歷史",
        "host 當大腦與對話介面,Fitness MCP 當背後的運動科學計算與安全判斷引擎。":
            "AI host 負責理解與表達;Pacevera 在使用者控制的資料邊界內負責計算、仲裁與追溯。",
        "AI 的教練能力:今天能不能練、該練哪裡、為什麼":
            "今天照不照原課表、具體改成什麼、為什麼、缺什麼",
        "Claude 和我們的 hosted MCP 看見完整的原始健康資料":
            "每個 AI host 或 Pacevera hosted service 各自保存完整健康歷史",
        "收最小化證據、算完不留 —— 這是產品的前提,不是附加的隱私功能。":
            "Local／private 保存 state;hosted 只做最小化 transient processing。兩者不能共用同一句隱私承諾。",
    },
    7: {
        "關鍵設計:我們什麼都不留":
            "關鍵設計:資料由使用者控制",
        "最小化 Evidence 進來,算出決策就丟掉 —— hosted 這端不寫任何一張表。":
            "Local／private 可保存可重算的 state、plan 與 decision history;hosted 不保存 raw Evidence。",
        "呼叫端送進來的": "使用者控制環境保存",
        "HRV 42ms · 睡眠 5h20m · 昨日 TSS 148":
            "Raw Evidence · provider token · longitudinal state · plan",
        "我們算完回傳的": "AI host 最小化取得",
        "Intervals → Zone 2 · ACWR 1.52 · 急性負荷偏快":
            "Decision from → to · reason · missing signals · versioned trace",
    },
    9: {
        "四層架構:原始資料停在界線之外":
            "四層架構:資料邊界與決策責任分開",
        "每一層只做一件事;最下面那一層,hosted MCP 不直接連。":
            "AI 可替換、Decision Engine 可重現、Evidence 與 token 留在 user-controlled boundary。",
        "Client Layer": "AI / Client Layer",
        "ChatGPT · Claude · Gemini · Cursor · Web · Mobile App · Apple Watch · Voice":
            "Claude · ChatGPT · Gemini · Web · Mobile · Team systems",
        "MCP Layer": "Pacevera Interface",
        "Fitness MCP —— 六個 tool,確定性計算,無狀態":
            "Desktop MCPB · private MCP · hosted remote · future REST / SDK",
        "Service Layer": "Decision Infrastructure",
        "Semantic Engine · Training Load · Decision Engine · Planning · Knowledge Graph":
            "Semantic State · Decision Engine · Rule Packages · Decision Graph",
        "Evidence 來源": "User-controlled Evidence",
        "Apple Health · Garmin · Strava —— 由 host 或使用者的 gateway 解析,不經過我們":
            "Local connectors · provider tokens · raw history · plans · outcomes",
    },
    18: {
        "OAuth 與同步在 Phase 2 的使用者環境執行,hosted 不持有 refresh token。":
            "Phase 2 在使用者控制環境完成授權、同步、正規化與撤銷;hosted 不持有 provider token。",
        "Webhook": "Revoke",
        "Fitbit": "Oura",
        "TrainingPeaks": "WHOOP",
        "Google Calendar": "Manual / CSV",
        "Notion": "Google Health",
        "新增一家資料源不用動上層;而在 Phase 1,這一層根本不在我們這邊。":
            "目前 6 家 parser;Oura／WHOOP 尚待真實去識別化 fixture。新增來源不改 Decision Engine。",
    },
    19: {
        "對外開放的六個 tool": "v0.5.0 規劃的十個 public tools",
        "這就是產品的能力清單 —— 任何支援 MCP 的客戶端都能直接呼叫。":
            "v0.4.2 已公開 6 個;v0.5.0 整合完成後新增 workout 與 coverage／trace／outcome。",
        "assess_fitness_state": "assess_fitness_state\nget_evidence_coverage",
        "從證據算出當前狀態\n": "狀態與訊號覆蓋\n",
        "decide_session": "decide_session\nexplain_decision",
        "今天的課表要不要改,改成什麼\n": "決策與版本化 trace\n",
        "generate_plan": "generate_plan\ngenerate_workout",
        "產生多週期計畫\n": "多週計畫與單次 workout\n",
        "commit_adjust_plan": "commit_adjust_plan\nsubmit_outcome",
        "兩階段提交,寫回計畫\n": "提交計畫與結果回報\n",
    },
    21: {
        "三種部署,一條資料界線": "三種部署,三句不同的隱私承諾",
        "同一套 domain package,差別只在 Evidence 在哪裡被處理。":
            "同一顆 Engine 與 Rule Packages;差別是 Evidence、token、state 與計算由誰控制。",
        "Hosted remote": "Local desktop",
        "短暫處理最小化 Evidence,永遠無狀態":
            "今天可用:本機 MCPB 與手動／匯出 Evidence",
        "Hosted 的界線": "User-controlled private",
        "不直連供應商、不持有 refresh token":
            "Phase 2:connectors、state、plan、decision 都在你的環境",
        "Local / private": "Controlled mobile",
        "connectors 與計算都在你的機器上跑":
            "Phase 3:pairing／tunnel 先連回 private engine",
        "Private data plane": "Hosted remote",
        "只有 user-controlled private 能做到原始資料不離開控制邊界":
            "只承諾最小化、transient、不留存;目前 Blocked",
    },
    30: {
        "現況與待做,一張表講完": "從已發行版本倒推下一階段",
        "已完成的附查證數字;待做的不寫沒有出處的驗收門檻。":
            "先關閉 v0.5.0,再用產品頁驗證市場,之後進入 private engine、mobile 與 team。",
        "資料標準化": "已發行 v0.4.2",
        "Apple Health／Garmin／Google／Strava／Oura／WHOOP／Oura／WHOOP":
            "Desktop MCPB · 6 public tools",
        "現況  6 家 parser,registry 6 家（8 種方言）":
            "Engine 1.6.0 · legacy Library 1.4.0",
        "確定性計算": "main 已完成",
        "ACWR／ATL／CTL／TSB 與分肌群疲勞":
            "trace／outcome／continuity／package lifecycle",
        "現況  470 tests、20 golden cases 全綠":
            "graph viewer · harness · regression gate",
        "六個對外 tool": "下一版 v0.5.0",
        "assess／decide／plan 三類":
            "10 tools · Engine 1.6.0",
        "現況  單次決策 0.443 ms":
            "base_rules 1.1.0 · planned",
        "Transport": "Phase 0",
        "stdio 與 Streamable HTTP": "整合 workout／防注入與 privacy",
        "現況  兩種都通": "pack · gates · publish",
        "OAuth Resource Server": "Phase 1",
        "簽章驗證、issuer、audience、scope": "pacevera.com 產品頁與真實 demo",
        "現況  private-development adapter；待做 production authorization server":
            "3–5 interviews · activation／retention",
        "遠端部署": "Phase 2",
        "HTTPS /mcp —— 手機場景的前提": "Local private engine MVP",
        "待做  桌機上架不需要它": "repository · connectors · continuity",
        "User-controlled private engine": "Phase 3",
        "connectors 與計算搬進使用者環境": "Controlled mobile · cross-host",
        "待做  local bundle／持久層／connector 權限": "hosted remote 維持 Blocked",
        "計畫的持久層": "Phase 4",
        "plan 與 planned workout 的表": "Team／Enterprise · REST／SDK",
        "待做  列在 Future Migrations": "governance · private deployment · pilot",
        "協定升級": "Remote GO gate",
        "2025-06-18 → 2026-07-28": "auth · HTTPS · redaction · privacy · E2E",
        "待做  走 dual-era,不直接切": "全部成立且有付費需求才 GO",
    },
    32: {
        "最重要的核心資產,\n是那個可重現的判斷,\n不是被保存下來的資料。":
            "最重要的核心資產,\n是連續 Evidence → 可追溯 Decision,\n而資料始終由使用者控制。",
        "下一步:先完成 P0 隱私契約與 P1 local private engine;remote 是後續 access channel,企業私有部署再後續":
            "下一步:先發布 v0.5.0 → pacevera.com 驗證 → local private engine;remote 仍是後續 access channel",
    },
}


def replace_text_preserving_runs(shape, old: str, new: str) -> bool:
    if not getattr(shape, "has_text_frame", False) or old not in shape.text:
        return False

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True

    # Preserve paragraph/run formatting when the replacement has the same
    # line structure (for example, the three-line takeaway slide).
    if shape.text == old:
        new_lines = new.split("\n")
        paragraphs = shape.text_frame.paragraphs
        if len(new_lines) == len(paragraphs) and all(p.runs for p in paragraphs):
            for paragraph, line in zip(paragraphs, new_lines):
                paragraph.runs[0].text = line
                for run in paragraph.runs[1:]:
                    run.text = ""
            return True
        shape.text = new
        return True

    return False


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()

    presentation = Presentation(args.source)
    missing: list[str] = []

    for slide_number, replacements in REPLACEMENTS.items():
        slide = presentation.slides[slide_number - 1]
        if all(
            any(
                getattr(shape, "has_text_frame", False) and new in shape.text
                for shape in slide.shapes
            )
            for new in replacements.values()
        ):
            continue
        for old, new in replacements.items():
            matches = [
                shape
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and old in shape.text
            ]
            if not matches and any(
                getattr(shape, "has_text_frame", False) and new in shape.text
                for shape in slide.shapes
            ):
                continue
            if len(matches) != 1:
                missing.append(
                    f"slide {slide_number}: expected one match for {old!r}, found {len(matches)}"
                )
                continue
            if not replace_text_preserving_runs(matches[0], old, new):
                missing.append(f"slide {slide_number}: could not replace {old!r}")

    if missing:
        raise SystemExit("\n".join(missing))

    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(args.output)


if __name__ == "__main__":
    main()
